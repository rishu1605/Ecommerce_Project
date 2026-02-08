# presentation_simple.py
# Run this script to create your PowerPoint file

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

print("🚀 Creating CoDeX SIC Mart Presentation...")

# Create a presentation
prs = Presentation()

# ==================== SLIDE 1: TITLE SLIDE ====================
slide1 = prs.slides.add_slide(prs.slide_layouts[0])

# Left side
left = slide1.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Team Name: CoDeX"
p.font.bold = True
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(0, 150, 255)

p = tf.add_paragraph()
p.text = "Project Title:"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(200, 200, 200)
p.space_after = Pt(12)

p = tf.add_paragraph()
p.text = "SIC Mart – A Multi-Vendor Cloud-Integrated\nE-Commerce Cart System"
p.font.bold = True
p.font.size = Pt(22)
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_after = Pt(24)

p = tf.add_paragraph()
p.text = "Team & Roles:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)

roles = [
    "Ramesh Kumar: Team Lead & Full Stack Architect",
    "Anmol Gautam: Logic & Financial Model",
    "Divyansh Sharma: Data Lead & DB Admin",
    "Akshat Dhiman: UI/UX & Frontend Designer",
    "Dhruv Paryag Sharma: Business Research"
]

for role in roles:
    p = tf.add_paragraph()
    p.text = f"• {role}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(220, 220, 220)
    p.level = 0
    p.space_before = Pt(6)

p = tf.add_paragraph()
p.text = "\nCourse: SIC Coding & Programming"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(18)

# ==================== SLIDE 2: PROBLEM STATEMENT ====================
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide2.shapes.title
title.text = "Problem Statement"

left = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Real-World Problem:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0, 200, 255)

p = tf.add_paragraph()
p.text = "• Vendors: Lack centralized tools for inventory & secure finance tracking."
p.font.size = Pt(14)
p.space_before = Pt(8)

p = tf.add_paragraph()
p.text = "• Buyers: Face data loss during fulfillment due to non-persistent shipping records."
p.font.size = Pt(14)
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "\nAffected Users:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(16)

p = tf.add_paragraph()
p.text = "Local micro-sellers & digital-native consumers."
p.font.size = Pt(14)
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "\nCore Issue:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 50, 50)

p = tf.add_paragraph()
p.text = 'Solving "Ghost Transactions" & "Data Decay"\nin multi-vendor e-commerce.'
p.font.size = Pt(15)
p.font.color.rgb = RGBColor(255, 100, 100)

# ==================== SLIDE 3: OBJECTIVE ====================
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide3.shapes.title
title.text = "Objective"

left = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Main Goals:"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0, 200, 255)

goals = [
    "1. Implement Atomic Financial Logic via a Virtual 'Sapphire Wallet'.",
    "2. Deploy Decoupled Media Storage using Cloudinary API.",
    "3. Establish Point-in-Time Data Persistence for order history.",
    "4. Provide a centralized Admin Bridge for oversight."
]

for goal in goals:
    p = tf.add_paragraph()
    p.text = goal
    p.font.size = Pt(14)
    p.space_before = Pt(8)
    p.level = 0

# ==================== SLIDE 4: PROPOSED SOLUTION ====================
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide4.shapes.title
title.text = "Proposed Solution"

left = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Approach:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)

p = tf.add_paragraph()
p.text = "Modular, three-tier ecosystem using Python & Streamlit."
p.font.size = Pt(14)
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "\nUnique Features:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(12)

features = [
    "• The Bridge: Centralized admin panel.",
    "• Hybrid Storage: SQLite (data) + Cloudinary (media).",
    "• SHA-256 Security: Enterprise-grade password hashing."
]

for feature in features:
    p = tf.add_paragraph()
    p.text = feature
    p.font.size = Pt(14)
    p.space_before = Pt(6)

# ==================== SLIDE 5: SYSTEM ARCHITECTURE ====================
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide5.shapes.title
title.text = "System Architecture"

left = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Data Flow:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)

p = tf.add_paragraph()
p.text = "Buyer Interface → Python Core (GST/Wallet Logic) → SQLite3 Persistence"
p.font.size = Pt(13)
p.space_before = Pt(8)

p = tf.add_paragraph()
p.text = "\nKey Principle:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(16)

p = tf.add_paragraph()
p.text = "Decoupled modules ensure seller inventory updates never disrupt buyer checkout sessions."
p.font.size = Pt(14)
p.space_before = Pt(6)

# ==================== SLIDE 6: TECH STACK ====================
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide6.shapes.title
title.text = "Tech Stack"

left = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

tech_items = [
    "• Language: Python 3.x",
    "• Frontend: Streamlit, Custom CSS/HTML",
    "• Backend: Pandas, Hashlib, UUID",
    "• Cloud API: Cloudinary",
    "• Database: SQLite3",
    "• Tools: VS Code, GitHub, secrets.toml"
]

for item in tech_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(14)
    p.space_before = Pt(6)

# ==================== SLIDE 7: CLOUD INTEGRATION ====================
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide7.shapes.title
title.text = "Cloud Integration"

left = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Dataset:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)

p = tf.add_paragraph()
p.text = "Custom relational schema in sic_mart.db"
p.font.size = Pt(14)
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "\nCloud-Decoupled Flow:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(12)

flow_steps = [
    "1. Input: Seller uploads image.",
    "2. Process: Cloudinary hosts & optimizes.",
    "3. Output: Secure HTTPS URL stored in DB."
]

for step in flow_steps:
    p = tf.add_paragraph()
    p.text = step
    p.font.size = Pt(14)
    p.space_before = Pt(6)

p = tf.add_paragraph()
p.text = "\nBenefit:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(100, 255, 100)
p.space_before = Pt(16)

p = tf.add_paragraph()
p.text = "80% smaller DB size; 0 server load for media."
p.font.size = Pt(14)
p.space_before = Pt(4)

# ==================== SLIDE 8: DATABASE SCHEMA ====================
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide8.shapes.title
title.text = "Database Schema"

left = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Relational Schema:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)

schema_items = [
    "• Users ↔ Addresses (1-to-Many)",
    "• Orders: Captures immutable price/address strings.",
    "• Triggers: Auto stock decrement in products."
]

for item in schema_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(14)
    p.space_before = Pt(6)

p = tf.add_paragraph()
p.text = "\nSample Columns:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "user_id, role, wallet_balance,\nimage_url (Cloudinary), order_status"
p.font.size = Pt(13)
p.space_before = Pt(4)

# ==================== SLIDE 9: UI & DATA FLOW ====================
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide9.shapes.title
title.text = "UI & Data Flow"

left = slide9.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Design:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)

p = tf.add_paragraph()
p.text = "Glassmorphism Dark Theme"
p.font.size = Pt(14)
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "\nAtomic Order Path:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(12)

order_steps = [
    "1. Cart Validation",
    "2. GST Calculation",
    "3. Wallet Deduction",
    "4. Record Commitment"
]

for step in order_steps:
    p = tf.add_paragraph()
    p.text = step
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(100, 200, 255)
    p.space_before = Pt(8)

p = tf.add_paragraph()
p.text = "\n[Presenter Note: Screenshots shown are from\nthe live Main.py running locally.]"
p.font.italic = True
p.font.size = Pt(11)
p.font.color.rgb = RGBColor(150, 150, 150)
p.space_before = Pt(16)

# ==================== SLIDE 10: RESULTS & PERFORMANCE ====================
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide10.shapes.title
title.text = "Results & Performance"

left = slide10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Outcome:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)

p = tf.add_paragraph()
p.text = "• 100% success in atomic transactions."
p.font.size = Pt(14)
p.space_before = Pt(6)

p = tf.add_paragraph()
p.text = "• 0 data conflicts in multi-address management."
p.font.size = Pt(14)
p.space_before = Pt(4)

p = tf.add_paragraph()
p.text = "\nKey Findings:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0, 200, 255)
p.space_before = Pt(12)

findings = [
    "• Cloudinary: Query times < 50ms.",
    "• SHA-256: 0-compromise security."
]

for finding in findings:
    p = tf.add_paragraph()
    p.text = finding
    p.font.size = Pt(14)
    p.space_before = Pt(6)

p = tf.add_paragraph()
p.text = "\nChallenge:"
p.font.bold = True
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 100, 100)
p.space_before = Pt(12)

p = tf.add_paragraph()
p.text = "Schema migrations & real-time state sync in Streamlit."
p.font.size = Pt(14)
p.space_before = Pt(4)

# ==================== SLIDE 11: FUTURE SCOPE ====================
slide11 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide11.shapes.title
title.text = "Future Scope"

left = slide11.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(5))
tf = left.text_frame
tf.clear()

future_items = [
    "• ML Integration: Buyer recommendation engine.",
    "• Logistics: Live GPS tracking module.",
    "• Real Payments: UPI/Card gateway bridges.",
    "• Scaling: Migrate from SQLite to PostgreSQL."
]

for item in future_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(14)
    p.space_before = Pt(8)

# ==================== SLIDE 12: THANK YOU ====================
slide12 = prs.slides.add_slide(prs.slide_layouts[0])

left = slide12.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(3))
tf = left.text_frame
tf.clear()

p = tf.add_paragraph()
p.text = "Thank You"
p.font.bold = True
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(0, 150, 255)
p.alignment = PP_ALIGN.LEFT

p = tf.add_paragraph()
p.text = "Questions?"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(200, 200, 200)
p.space_before = Pt(24)

p = tf.add_paragraph()
p.text = "\nSamsung Innovation Campus"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(100, 100, 255)
p.space_before = Pt(36)

# ==================== SAVE THE FILE ====================
file_name = "CoDeX_SIC_Mart_Presentation.pptx"
prs.save(file_name)

full_path = os.path.abspath(file_name)
print(f"✅ Presentation created successfully!")
print(f"📁 File saved as: {full_path}")
print(f"\n📋 To open: Double-click the file or open in Microsoft PowerPoint")
print(f"📍 Location: {os.path.dirname(full_path)}")